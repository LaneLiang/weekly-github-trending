"""Presentation workflow — generates academic group-meeting PPTX following nature-paper2ppt design rules.

Layout philosophy: varied compositions, evidence-first, no repetitive AI-template cards.
"""
from datetime import date
from pathlib import Path

from lanes_ceo.contracts import Artifact, CriticReview, Job

# ── Design tokens (Nature-style academic) ──
_COLOR_BG = (0xFF, 0xFF, 0xFF)        # white background
_COLOR_PRIMARY = (0x1A, 0x1A, 0x2E)    # dark navy — titles
_COLOR_BODY = (0x33, 0x33, 0x33)       # dark gray — body text
_COLOR_ACCENT = (0x2C, 0x5F, 0x8A)     # academic blue — accents
_COLOR_MUTED = (0x88, 0x88, 0x88)      # gray — captions/sources
_COLOR_LIGHT_BG = (0xF5, 0xF7, 0xFA)  # light blue-gray — accent strips
_COLOR_DIVIDER = (0xDD, 0xDD, 0xDD)    # light divider

_FONT_TITLE_CN = "微软雅黑"
_FONT_BODY_CN = "微软雅黑"
_FONT_EN = "Arial"

_TITLE_SIZE = 28          # pt — slide title
_SUBTITLE_SIZE = 14       # pt — subtitle/metadata
_BODY_SIZE = 14           # pt — bullet body
_SMALL_SIZE = 9           # pt — source/caption
_METRIC_SIZE = 36         # pt — large metric numbers

_SLIDE_W = 13350400       # 16:9 EMU
_SLIDE_H = 7512000

PRESENTATION_STRUCTURE = [
    ("封面", "汇报标题 / 姓名 / 日期"),
    ("目录", "汇报大纲"),
    ("研究背景与意义", "课题来源、学术价值、应用前景"),
    ("本周核心进展", "2-3 项关键进展，每项配数据/图表支撑"),
    ("实验结果与分析", "实验设计、数据展示、结果讨论"),
    ("难点与挑战", "遇到的技术瓶颈和理论困难"),
    ("方案对比与选型", "候选方案对比、决策依据"),
    ("下周工作计划", "具体任务、预期产出、时间节点"),
    ("需要讨论的问题", "向导师和团队请教的问题"),
    ("参考文献与致谢", "引用文献、合作致谢"),
]

TARGET_SLIDE_COUNT = 15
MIN_SLIDES = 12
MAX_SLIDES = 18
CORE_SLIDE_MIN_RATIO = 0.30


class PresentationWorkflow:
    role_group = "presentation"
    actor_name = "presentation-actor"
    critic_name = "presentation-critic"

    def run_actor(self, job: Job) -> Artifact:
        message = job.input.get("message", "生成组会汇报PPT")
        today = date.today()
        week_label = f"{today.year}W{today.isocalendar()[1]}"

        from lanes_ceo.workflows.utils import get_artifact_dir, llm_chat

        prompt = _build_ppt_prompt(message)
        llm_response = llm_chat(
            "你是一名博士研究生，需准备30分钟组会PPT汇报。"
            "为每页输出：标题行（结论式，如'本周完成3项核心实验设计'）"
            "之后是3-5个中文要点，每条要点不超过20字。",
            prompt,
        ) or ""

        artifact_dir = get_artifact_dir("presentation")
        ts = today.strftime("%H%M%S")
        pptx_path = artifact_dir / f"presentation-{week_label}-{ts}.pptx"

        slide_count = _write_pptx(pptx_path, message, llm_response)

        return Artifact(
            artifact_id=f"artifact-{job.job_id}",
            job_id=job.job_id,
            artifact_type="presentation",
            summary=f"PPT 汇报 {week_label}: {message} ({slide_count} 页)",
            artifact_paths=[str(pptx_path)],
            sources=["llm-generated", "user-input"],
            risks=["LLM 内容需人工核实", "图表需手动插入"],
            user_confirmations=[
                "请确认幻灯片内容和数据准确性",
                "请确认汇报时长是否适配 (~30分钟)",
            ],
        )

    def run_critic(self, job: Job, artifact: Artifact) -> CriticReview:
        issues: list[str] = []
        score = 85

        if artifact.artifact_paths:
            path = Path(artifact.artifact_paths[0])
            if not path.exists():
                issues.append("生成的 pptx 文件不存在")
                score -= 30

        slide_count = _extract_slide_count(artifact.summary)
        if slide_count is not None:
            if slide_count < MIN_SLIDES:
                issues.append(f"页数不足 (当前 {slide_count}, 最少 {MIN_SLIDES})")
                score -= 10
            elif slide_count > MAX_SLIDES:
                issues.append(f"页数过多 (当前 {slide_count}, 最多 {MAX_SLIDES})")
                score -= 10

        structure_ok = _check_structure(artifact.summary)
        if not structure_ok:
            issues.append("PPT 结构不完整，缺少关键章节")

        approved = len(issues) == 0
        return CriticReview(
            review_id=f"review-{job.job_id}",
            job_id=job.job_id,
            review_result="approved" if approved else "returned",
            score=max(score, 0),
            issues=issues,
            approved=approved,
            return_to_actor=not approved,
            handoff_note="PPT 质量合格" if approved else "请修正后重新提交",
        )


def _build_ppt_prompt(message: str) -> str:
    structure_text = "\n".join(
        f"{i + 1}. {s[0]} — {s[1]}" for i, s in enumerate(PRESENTATION_STRUCTURE)
    )
    return (
        f"请根据以下主题生成组会汇报PPT的完整大纲，共约{TARGET_SLIDE_COUNT}页：\n\n"
        f"主题：{message}\n\n"
        f"要求结构：\n{structure_text}\n\n"
        f"要点：每页给出标题和3-5个要点，核心成果页占比≥{int(CORE_SLIDE_MIN_RATIO * 100)}%。"
    )


# ═══════════════════════════════════════════════════════════════
# PPTX builder — nature-paper2ppt style
# ═══════════════════════════════════════════════════════════════

def _write_pptx(path: Path, message: str, llm_response: str) -> int:
    from pptx import Presentation
    from pptx.util import Pt, Inches, Emu, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    # Use blank layout as base; we build everything via text boxes
    blank_layout = prs.slide_layouts[6]  # blank

    title_slide = prs.slides.add_slide(blank_layout)
    _build_title_slide(title_slide, message)

    if llm_response and not llm_response.startswith("[LLM"):
        slides_data = _parse_llm_slides(llm_response)
    else:
        slides_data = [
            (title, f"（{desc} — 待 LLM 填充）")
            for title, desc in PRESENTATION_STRUCTURE
        ]

    # Track layout pattern usage for anti-template check
    _last_layout = None
    for idx, (title, body) in enumerate(slides_data):
        # Rotate through layout types to avoid repetition
        layout_type = idx % 5
        slide = prs.slides.add_slide(blank_layout)
        _build_content_slide(slide, title, body, idx + 2, layout_type)

    prs.save(str(path))
    return len(prs.slides)


def _build_title_slide(slide, title_text: str) -> None:
    """Nature-style cover: large title, small metadata, clean."""
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Dark accent bar at top
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(0), Emu(0), _SLIDE_W, Emu(600000),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*_COLOR_PRIMARY)
    bar.line.fill.background()

    # Title text box
    txBox = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.2), Inches(11.33), Inches(2.0),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*_COLOR_PRIMARY)
    p.font.name = _FONT_TITLE_CN
    p.alignment = PP_ALIGN.LEFT

    # Subtitle: date + context
    txSub = slide.shapes.add_textbox(
        Inches(1.0), Inches(4.4), Inches(11.33), Inches(0.8),
    )
    tf2 = txSub.text_frame
    p2 = tf2.paragraphs[0]
    today = date.today()
    week_label = f"{today.year}W{today.isocalendar()[1]}"
    p2.text = f"组会汇报  |  {week_label}  |  东南大学"
    p2.font.size = Pt(_SUBTITLE_SIZE)
    p2.font.color.rgb = RGBColor(*_COLOR_MUTED)
    p2.font.name = _FONT_BODY_CN
    p2.alignment = PP_ALIGN.LEFT

    # Thin accent line under subtitle
    line = slide.shapes.add_shape(
        1,
        Inches(1.0), Inches(5.35), Inches(2.5), Emu(36000),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*_COLOR_ACCENT)
    line.line.fill.background()

    # Bottom metadata bar
    bot_bar = slide.shapes.add_shape(
        1, Emu(0), _SLIDE_H - Emu(450000), _SLIDE_W, Emu(450000),
    )
    bot_bar.fill.solid()
    bot_bar.fill.fore_color.rgb = RGBColor(*_COLOR_PRIMARY)
    bot_bar.line.fill.background()


def _build_content_slide(slide, title: str, body: str, slide_num: int, layout_type: int) -> None:
    """Build a content slide with varied composition based on layout_type."""
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    # All slides share: title bar at top, page number at bottom right
    _add_slide_title_bar(slide, title)
    _add_page_number(slide, slide_num)

    body_lines = [l.strip().lstrip("-•· ") for l in body.split("\n") if l.strip()]

    if layout_type == 0:
        _layout_hero_text(slide, body_lines)
    elif layout_type == 1:
        _layout_process_wide(slide, body_lines)
    elif layout_type == 2:
        _layout_claim_cards(slide, body_lines)
    elif layout_type == 3:
        _layout_comparison(slide, body_lines)
    else:
        _layout_discussion(slide, body_lines)


def _add_slide_title_bar(slide, title: str) -> None:
    """Dark title bar spanning the top of the slide."""
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor

    # Title background strip
    bar = slide.shapes.add_shape(
        1, Emu(0), Emu(0), _SLIDE_W, Emu(950000),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*_COLOR_PRIMARY)
    bar.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.22), Inches(11.57), Inches(0.6),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(_TITLE_SIZE)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = _FONT_TITLE_CN


def _add_page_number(slide, num: int) -> None:
    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    txBox = slide.shapes.add_textbox(
        Inches(11.5), Inches(6.85), Inches(1.5), Inches(0.4),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(_SMALL_SIZE)
    p.font.color.rgb = RGBColor(*_COLOR_MUTED)
    p.font.name = _FONT_EN
    p.alignment = PP_ALIGN.RIGHT


# ── Layout variants (anti-template: 5 different compositions) ──


def _layout_hero_text(slide, lines: list[str]) -> None:
    """Large claim-led layout: one big statement + supporting fragments."""
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if not lines:
        return
    # Main claim — large
    txMain = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.6), Inches(11.33), Inches(1.8),
    )
    tf = txMain.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = lines[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*_COLOR_PRIMARY)
    p.font.name = _FONT_TITLE_CN

    # Accent line
    line = slide.shapes.add_shape(
        1, Inches(1.0), Inches(3.55), Inches(2.0), Emu(24000),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*_COLOR_ACCENT)
    line.line.fill.background()

    # Supporting bullets
    if len(lines) > 1:
        txBody = slide.shapes.add_textbox(
            Inches(1.0), Inches(3.8), Inches(11.33), Inches(2.5),
        )
        tf2 = txBody.text_frame
        tf2.word_wrap = True
        for i, line_text in enumerate(lines[1:5]):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = f"  {line_text}"
            p2.font.size = Pt(_BODY_SIZE)
            p2.font.color.rgb = RGBColor(*_COLOR_BODY)
            p2.font.name = _FONT_BODY_CN
            p2.space_after = Pt(8)


def _layout_process_wide(slide, lines: list[str]) -> None:
    """Full-width process/workflow layout with stage labels."""
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    n = min(len(lines), 4)
    if n == 0:
        return

    box_w = Inches(2.6)
    box_h = Inches(2.4)
    start_x = Inches(0.7)
    gap = Inches(0.25)
    y = Inches(1.9)

    colors = [
        RGBColor(*_COLOR_PRIMARY),
        RGBColor(*_COLOR_ACCENT),
        RGBColor(0x3D, 0x7A, 0xBF),
        RGBColor(0x5A, 0x9E, 0xD4),
    ]

    for i in range(n):
        x = int(start_x + i * (box_w + gap))

        # Card background
        card = slide.shapes.add_shape(
            1, x, y, box_w, box_h,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = colors[i]
        card.line.fill.background()

        # Step number
        num_box = slide.shapes.add_textbox(
            x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.4),
        )
        np = num_box.text_frame.paragraphs[0]
        np.text = f"0{i + 1}"
        np.font.size = Pt(11)
        np.font.bold = True
        np.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        np.font.name = _FONT_EN

        # Step text
        text_box = slide.shapes.add_textbox(
            x + Inches(0.2), y + Inches(0.55), Inches(box_w - Inches(0.4)), Inches(1.6),
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        tp = tf.paragraphs[0]
        tp.text = lines[i]
        tp.font.size = Pt(_BODY_SIZE)
        tp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tp.font.name = _FONT_BODY_CN


def _layout_claim_cards(slide, lines: list[str]) -> None:
    """3 compact claim cards in a row, each with one strong sentence."""
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    n = min(len(lines), 3)
    if n == 0:
        return

    card_w = Inches(3.5)
    card_h = Inches(1.8)
    y = Inches(2.0)
    total_w = n * card_w + (n - 1) * Inches(0.3)
    start_x = int((_SLIDE_W - total_w) / 2)

    for i in range(n):
        x = int(start_x + i * (card_w + Inches(0.3)))

        # Card background
        card = slide.shapes.add_shape(1, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*_COLOR_LIGHT_BG)
        card.line.color.rgb = RGBColor(*_COLOR_DIVIDER)
        card.line.width = Pt(0.5)

        # Card number badge
        badge = slide.shapes.add_shape(
            1, x + Inches(0.25), y + Inches(0.2), Inches(0.45), Inches(0.45),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(*_COLOR_ACCENT)
        badge.line.fill.background()
        bp = badge.text_frame.paragraphs[0]
        bp.text = str(i + 1)
        bp.font.size = Pt(14)
        bp.font.bold = True
        bp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        bp.font.name = _FONT_EN
        bp.alignment = PP_ALIGN.CENTER

        # Card text
        txBox = slide.shapes.add_textbox(
            x + Inches(0.3), y + Inches(0.8), card_w - Inches(0.6), Inches(0.85),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tp = tf.paragraphs[0]
        tp.text = lines[i]
        tp.font.size = Pt(_BODY_SIZE)
        tp.font.color.rgb = RGBColor(*_COLOR_BODY)
        tp.font.name = _FONT_BODY_CN


def _layout_comparison(slide, lines: list[str]) -> None:
    """Left/right comparison with a conclusion strip at bottom."""
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if not lines:
        return

    # Split into two columns
    mid = (len(lines) + 1) // 2
    left_lines = lines[:mid]
    right_lines = lines[mid:]

    col_y = Inches(1.8)
    col_w = Inches(5.5)

    for col_idx, col_lines in [(0, left_lines), (1, right_lines)]:
        x = Inches(0.8) if col_idx == 0 else Inches(7.0)

        for j, line_text in enumerate(col_lines[:4]):
            txBox = slide.shapes.add_textbox(
                x, int(col_y + j * Inches(1.1)), col_w, Inches(0.9),
            )
            # Small bullet indicator
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"  {line_text}"
            p.font.size = Pt(_BODY_SIZE)
            p.font.color.rgb = RGBColor(*_COLOR_BODY)
            p.font.name = _FONT_BODY_CN

    # Bottom conclusion strip
    strip = slide.shapes.add_shape(
        1, Inches(0.8), Inches(6.1), Inches(11.57), Inches(0.6),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(*_COLOR_PRIMARY)
    strip.line.fill.background()
    sp = strip.text_frame.paragraphs[0]
    sp.text = "  结论" if left_lines else ""
    sp.font.size = Pt(11)
    sp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sp.font.name = _FONT_BODY_CN


def _layout_discussion(slide, lines: list[str]) -> None:
    """Open discussion layout: questions with generous spacing."""
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor

    if not lines:
        return

    for i, line_text in enumerate(lines[:5]):
        y = Inches(1.8) + i * Inches(1.0)
        txBox = slide.shapes.add_textbox(
            Inches(1.5), y, Inches(10.33), Inches(0.8),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"  {line_text}"
        p.font.size = Pt(_BODY_SIZE + 2)
        p.font.color.rgb = RGBColor(*_COLOR_BODY)
        p.font.name = _FONT_BODY_CN

        # Thin left accent mark
        mark = slide.shapes.add_shape(
            1,
            Inches(1.1), y + Inches(0.08),
            Emu(72000), Emu(500000),
        )
        mark.fill.solid()
        mark.fill.fore_color.rgb = RGBColor(*_COLOR_ACCENT)
        mark.line.fill.background()


def _parse_llm_slides(text: str) -> list[tuple[str, str]]:
    """Parse LLM output into (title, body) pairs."""
    result = []
    blocks = text.split("\n\n")
    for block in blocks:
        stripped = block.strip()
        if not stripped:
            continue
        lines = stripped.split("\n", 1)
        title = lines[0].lstrip("0123456789. )-【】#")
        body = lines[1] if len(lines) > 1 else ""
        result.append((title, body))
    return result


def _extract_slide_count(summary: str) -> int | None:
    import re
    match = re.search(r"(\d+)\s*页", summary)
    return int(match.group(1)) if match else None


def _check_structure(summary: str) -> bool:
    required_sections = ["背景", "进展", "实验", "计划"]
    return any(s in summary for s in required_sections)
